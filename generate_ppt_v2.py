from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

OUTPUT_PPT = r"C:\Users\ADVAN\brightdor\BrightDor_Documentation.pptx"
OUTPUT_PPT_DOCS = r"C:\Users\ADVAN\brightdor\docs\BrightDor-Dokumentasi-PPT1.pptx"
SS_DIR = r"C:\Users\ADVAN\brightdor\screenshots"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

CREAM = RGBColor(0xC4, 0xA5, 0x74)
CREAM_LIGHT = RGBColor(0xF5, 0xF0, 0xE8)
CREAM_BG = RGBColor(0xFA, 0xF7, 0xF2)
BLACK = RGBColor(0x14, 0x14, 0x14)
DARK_GRAY = RGBColor(0x2A, 0x2A, 0x2A)
GRAY = RGBColor(0x52, 0x52, 0x5B)
LIGHT_GRAY = RGBColor(0xA1, 0xA1, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x16, 0xA3, 0x4A)
WARNING = RGBColor(0xEA, 0xB3, 0x08)
DANGER = RGBColor(0xDC, 0x26, 0x26)
INFO = RGBColor(0x25, 0x63, 0xEB)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)

FONT_MAIN = "Calibri"

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or WHITE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or WHITE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, text, left, top, width, height, font_size=Pt(14), color=BLACK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT_MAIN
    p.alignment = alignment
    return txBox

def add_multi_text(slide, items, left, top, width, font_size=Pt(12), color=DARK_GRAY, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "  •  " + item
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = FONT_MAIN
        p.space_after = spacing
    return txBox

def add_line(slide, left, top, width, color=CREAM):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False

def add_top_bar(slide):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=CREAM)

def add_bottom_bar(slide):
    add_rect(slide, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), fill_color=CREAM)

def add_slide_number(slide, num, total=16):
    add_text(slide, f"{num} / {total}", Inches(12.2), Inches(7.05), Inches(1), Inches(0.35),
             font_size=Pt(10), color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, title, subtitle=""):
    add_top_bar(slide)
    add_text(slide, title, Inches(0.8), Inches(0.35), Inches(10), Inches(0.55),
             font_size=Pt(28), color=BLACK, bold=True)
    add_line(slide, Inches(0.8), Inches(0.95), Inches(1.5), CREAM)
    if subtitle:
        add_text(slide, subtitle, Inches(0.8), Inches(1.05), Inches(10), Inches(0.4),
                 font_size=Pt(13), color=GRAY)

def add_screenshot(slide, filename, left, top, max_width, max_height):
    path = os.path.join(SS_DIR, filename)
    if not os.path.exists(path):
        add_text(slide, f"[Screenshot: {filename}]", left, top, max_width, Inches(0.4), color=DANGER)
        return
    img = Image.open(path)
    img_w, img_h = img.size
    ratio = min(max_width / img_w, max_height / img_h)
    new_w = int(img_w * ratio)
    new_h = int(img_h * ratio)
    slide.shapes.add_picture(path, left, top, Emu(new_w), Emu(new_h))

def add_stat_card(slide, left, top, label, value, accent_color=CREAM):
    card = add_rounded_rect(slide, left, top, Inches(2.2), Inches(1.15), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5))
    add_rect(slide, left, top + Inches(0.95), Inches(2.2), Inches(0.06), fill_color=accent_color)
    add_text(slide, value, left + Inches(0.15), top + Inches(0.12), Inches(1.9), Inches(0.45),
             font_size=Pt(22), color=BLACK, bold=True)
    add_text(slide, label, left + Inches(0.15), top + Inches(0.55), Inches(1.9), Inches(0.35),
             font_size=Pt(11), color=GRAY)


# ==================== SLIDE 1: COVER ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_rect(slide, Inches(0), Inches(0), Inches(0.12), prs.slide_height, fill_color=CREAM)
add_rect(slide, Inches(0), Inches(3.2), prs.slide_width, Inches(0.005), fill_color=CREAM)

add_text(slide, "BRIGHTDOR", Inches(1.2), Inches(1.8), Inches(11), Inches(1),
         font_size=Pt(52), color=CREAM, bold=True)
add_text(slide, "Documentation Project", Inches(1.2), Inches(2.55), Inches(11), Inches(0.6),
         font_size=Pt(22), color=WHITE)
add_text(slide, "Marketplace Jasa Pernikahan Premium & Undangan Digital", Inches(1.2), Inches(3.5), Inches(11), Inches(0.5),
         font_size=Pt(16), color=LIGHT_GRAY)

info_items = [
    "Laravel 11  |  Filament v5  |  SQLite  |  37 Tables (sync 1 Sep 2026)",
    "Teknologi Informasi — 2026",
]
y = Inches(5.2)
for item in info_items:
    add_text(slide, item, Inches(1.2), y, Inches(11), Inches(0.35), font_size=Pt(13), color=LIGHT_GRAY)
    y += Inches(0.35)

add_rect(slide, Inches(11.5), Inches(0), Inches(1.833), prs.slide_height, fill_color=CREAM)
add_text(slide, "2026", Inches(11.6), Inches(3.2), Inches(1.6), Inches(0.6),
         font_size=Pt(28), color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)


# ==================== SLIDE 2: NAMA APLIKASI ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "BrightDor", "Marketplace Jasa Pernikahan Premium & Undangan Digital")
add_bottom_bar(slide)
add_slide_number(slide, 2)

add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.8), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5))
add_text(slide, "Tentang BrightDor", Inches(1.1), Inches(1.6), Inches(5), Inches(0.4),
         font_size=Pt(18), color=BLACK, bold=True)
desc = (
    "BrightDor adalah platform digital yang menghubungkan pasangan pengantin (couple) "
    "dengan vendor jasa pernikahan premium di Indonesia. Platform ini menyediakan "
    "marketplace lengkap untuk layanan pernikahan beserta fitur undangan digital "
    "interaktif dengan sistem RSVP online."
)
add_text(slide, desc, Inches(1.1), Inches(2.05), Inches(11.2), Inches(1.1),
         font_size=Pt(13), color=GRAY)

features = [
    ("Vendor Management", "Approval, kategori, & verifikasi vendor"),
    ("Marketplace Jasa", "Katalog layanan pernikahan premium"),
    ("Booking & Order", "Sistem pemesanan & tracking status"),
    ("Undangan Digital", "Template premium + RSVP online"),
    ("Sistem Keuangan", "Transaksi, komisi, & payout vendor"),
    ("Multi Bahasa", "8 bahasa termasuk Indonesia & Arab"),
]
for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    fx = Inches(0.8) + Inches(col * 4.0)
    fy = Inches(3.6) + Inches(row * 1.6)
    fcard = add_rounded_rect(slide, fx, fy, Inches(3.7), Inches(1.25), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5))
    add_rect(slide, fx, fy, Inches(3.7), Inches(0.05), fill_color=CREAM)
    add_text(slide, title, fx + Inches(0.15), fy + Inches(0.15), Inches(3.4), Inches(0.35),
             font_size=Pt(13), color=BLACK, bold=True)
    add_text(slide, desc, fx + Inches(0.15), fy + Inches(0.55), Inches(3.4), Inches(0.55),
             font_size=Pt(11), color=GRAY)


# ==================== SLIDE 3: HOME & NAVBAR — TERBARU 1 Sep 2026 (menggantikan navbar sempit) ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Pengalaman Pengunjung — Navigasi & Kategori (1 Sep 2026)", "Satu baris .bd-catnav + ikon Heroicons konsisten + modal booking — menggantikan flex-wrap sempit")
add_bottom_bar(slide)
add_slide_number(slide, 3)

# Two home screenshots side by side (01_home_top + 02_home_scrolled) plus navbar_v2 glass
add_screenshot(slide, "01_home_top.png", Inches(0.3), Inches(1.45), Inches(6.2), Inches(3.4))
add_screenshot(slide, "02_home_scrolled.png", Inches(6.8), Inches(1.45), Inches(6.2), Inches(3.4))
add_screenshot(slide, r"navbar_v2\02b_desktop_glass_clean.png", Inches(0.3), Inches(5.0), Inches(6.2), Inches(2.2))
add_screenshot(slide, r"navbar_v2\04e_desktop_dropdown_full.png", Inches(6.8), Inches(5.0), Inches(6.2), Inches(2.2))

# caption box
cap = add_rounded_rect(slide, Inches(0.3), Inches(7.25), Inches(12.7), Inches(0.35), fill_color=WHITE, border_color=RGBColor(0xE5,0xE5,0xE5))
add_text(slide, "  Kiri: home top (hero + kategori)  |  Kanan: scrolled (navbar glass)  |  Bawah: dropdown glass + Ajukan Penawaran modal — file: navigation.blade.php / category-icon.blade.php / vendors/show.blade.php", Inches(0.4), Inches(7.28), Inches(12.5), Inches(0.3), font_size=Pt(8), color=GRAY)

# ==================== SLIDE 3b: KATEGORI FOTOGRAFER — CARD VENDOR DENGAN GAMBAR PORTFOLIO (menggantikan inisial) ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Kategori Fotografer — Card Vendor dengan Gambar Portfolio", "Menggantikan inisial dengan cover image via Spatie Media Library (services.cover + vendors.portfolio)")
add_bottom_bar(slide)
add_slide_number(slide, 4)

add_screenshot(slide, "02b_category_fotografer.png", Inches(0.5), Inches(1.45), Inches(8.5), Inches(5.8))

expl = add_rounded_rect(slide, Inches(9.3), Inches(1.45), Inches(3.7), Inches(5.8), fill_color=WHITE, border_color=RGBColor(0xE5,0xE5,0xE5))
add_rect(slide, Inches(9.3), Inches(1.45), Inches(3.7), Inches(0.4), fill_color=CREAM)
add_text(slide, "  Perbaikan Card Vendor", Inches(9.5), Inches(1.48), Inches(3.3), Inches(0.35), font_size=Pt(11), color=BLACK, bold=True)
add_multi_text(slide, [
    "Sebelum: card hanya inisial huruf, tanpa gambar",
    "Sesudah: cover image dari services.cover (Spatie singleFile)",
    "Fallback ke vendor.portfolio jika cover kosong",
    "Kategori satu baris .bd-catnav + ikon Heroicons konsisten",
    "Storage: public disk, conversions, responsive",
    "Screenshot: 02b_category_fotografer.png (1 Sep, fullpage)",
], Inches(9.5), Inches(2.0), Inches(3.3), Pt(10), DARK_GRAY, Pt(7))

# ==================== SLIDE 4: DASHBOARD (REAL SCREENSHOT — UPDATED 1 Sep 2026) ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Dashboard Admin — Terbaru (1 Sep 2026)", "Top navigation light (cream #C4A574) • global search Cmd+K • warna putih / cream / hitam — sync dengan perbaikan navbar")
add_bottom_bar(slide)
add_slide_number(slide, 5)

# Use latest screenshot dated 1 Sep 2026 (04_admin_dashboard.png = 501KB) — previously dashboard.png (461KB, 31 Aug, dark sidebar)
add_screenshot(slide, "04_admin_dashboard.png", Inches(0.5), Inches(1.45), Inches(8.5), Inches(5.8))

explanation_box = add_rounded_rect(slide, Inches(9.3), Inches(1.45), Inches(3.7), Inches(5.8), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5))
add_rect(slide, Inches(9.3), Inches(1.45), Inches(3.7), Inches(0.4), fill_color=BLACK)
add_text(slide, "  Fitur Dashboard", Inches(9.5), Inches(1.48), Inches(3.3), Inches(0.35),
         font_size=Pt(12), color=CREAM, bold=True)
add_multi_text(slide, [
    "5 kartu statistik: Vendors, Bookings, Revenue, Couples, Undangan",
    "Revenue chart 6 bulan (line chart) + doughnut bookings/kategori",
    "Quick Actions + Recent Activities (vendor pending & booking terbaru)",
    "Top Navigation LIGHT: tanpa sidebar, lega, krem/hitam/putih (fixed vs dark old)",
    "Global search Cmd+K, notif, profil — sinkron perbaikan 1 Sep",
    "Screenshot: 04_admin_dashboard.png (1 Sep 2026, 1440x900 @2x)",
], Inches(9.5), Inches(2.0), Inches(3.3), Pt(11), DARK_GRAY, Pt(8))


# ==================== SLIDE 4: MANAJEMEN VENDOR ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Mockup: Manajemen Vendor", "Approval, kategori, verifikasi, dan pengelolaan vendor pernikahan")
add_bottom_bar(slide)
add_slide_number(slide, 6)

add_screenshot(slide, "vendors.png", Inches(0.5), Inches(1.45), Inches(8.5), Inches(5.8))

explanation_box = add_rounded_rect(slide, Inches(9.3), Inches(1.45), Inches(3.7), Inches(5.8), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5))
add_rect(slide, Inches(9.3), Inches(1.45), Inches(3.7), Inches(0.4), fill_color=CREAM)
add_text(slide, "  Fitur Vendor", Inches(9.5), Inches(1.48), Inches(3.3), Inches(0.35),
         font_size=Pt(12), color=BLACK, bold=True)
add_multi_text(slide, [
    "Tabel dengan kolom: nama, pemilik, kategori, kota, rating, status",
    "Badge status: Pending (kuning), Approved (hijau), Rejected (merah)",
    "Aksi: Approve, Reject, Suspend dengan konfirmasi",
    "Filter & search untuk pencarian vendor",
    "Vendor Categories: 10 kategori (Venue, Catering, Dekorasi, dll)",
    "Form profil usaha: nama, deskripsi, lokasi, kontak, bank",
    "Rating & verifikasi vendor",
    "Badge warning untuk vendor pending di navigasi",
], Inches(9.5), Inches(2.0), Inches(3.3), Pt(11), DARK_GRAY, Pt(6))


# ==================== SLIDE 5: BOOKING/ORDER ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Mockup: Booking & Order", "Sistem pemesanan layanan vendor dengan tracking status real-time")
add_bottom_bar(slide)
add_slide_number(slide, 7)

add_screenshot(slide, "bookings.png", Inches(0.5), Inches(1.45), Inches(8.5), Inches(5.8))

explanation_box = add_rounded_rect(slide, Inches(9.3), Inches(1.45), Inches(3.7), Inches(5.8), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5))
add_rect(slide, Inches(9.3), Inches(1.45), Inches(3.7), Inches(0.4), fill_color=INFO)
add_text(slide, "  Fitur Booking", Inches(9.5), Inches(1.48), Inches(3.3), Inches(0.35),
         font_size=Pt(12), color=WHITE, bold=True)
add_multi_text(slide, [
    "Status flow: Pending → Confirmed → On Progress → Completed",
    "Opsi: Cancel / Refund dengan alasan",
    "Data acara: tanggal, waktu, lokasi, jumlah tamu",
    "Komponen keuangan: subtotal, diskon, admin fee, komisi, total",
    "Catatan customer & catatan admin",
    "Auto-generated booking code: BD-{8 karakter}",
    "Badge counter di navigasi untuk booking pending",
    "Search & filter berdasarkan status dan tanggal",
], Inches(9.5), Inches(2.0), Inches(3.3), Pt(11), DARK_GRAY, Pt(6))


# ==================== SLIDE 6: UNDANGAN DIGITAL ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Mockup: Undangan Digital", "Template premium undangan digital + RSVP online interaktif")
add_bottom_bar(slide)
add_slide_number(slide, 8)

add_screenshot(slide, "invitation-templates.png", Inches(0.5), Inches(1.45), Inches(8.5), Inches(5.8))

explanation_box = add_rounded_rect(slide, Inches(9.3), Inches(1.45), Inches(3.7), Inches(5.8), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5))
add_rect(slide, Inches(9.3), Inches(1.45), Inches(3.7), Inches(0.4), fill_color=PURPLE)
add_text(slide, "  Fitur Undangan Digital", Inches(9.5), Inches(1.48), Inches(3.3), Inches(0.35),
         font_size=Pt(12), color=WHITE, bold=True)
add_multi_text(slide, [
    "Template: Elegant Gold, Garden Romance, Modern Minimalis, dll",
    "Harga per template: Rp 150.000 - Rp 300.000",
    "Kategori: Elegant, Modern, Traditional",
    "Fitur: RSVP Online, Love Story, Gallery, Countdown, Music",
    "Status Premium & Featured pada template",
    "Order: Nama pengantin, tanggal, venue, subdomain",
    "Custom domain support untuk undangan",
    "Invitation live dengan views & RSVP counters",
    "3 tahap: Pilih Template → Bayar → Publish",
], Inches(9.5), Inches(2.0), Inches(3.3), Pt(11), DARK_GRAY, Pt(5))


# ==================== SLIDE 7: KEUANGAN ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Mockup: Keuangan", "Transaksi, komisi vendor, payout, dan pengaturan komisi")
add_bottom_bar(slide)
add_slide_number(slide, 9)

add_screenshot(slide, "transactions.png", Inches(0.5), Inches(1.45), Inches(8.5), Inches(5.8))

explanation_box = add_rounded_rect(slide, Inches(9.3), Inches(1.45), Inches(3.7), Inches(5.8), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5))
add_rect(slide, Inches(9.3), Inches(1.45), Inches(3.7), Inches(0.4), fill_color=SUCCESS)
add_text(slide, "  Fitur Keuangan", Inches(9.5), Inches(1.48), Inches(3.3), Inches(0.35),
         font_size=Pt(12), color=WHITE, bold=True)
add_multi_text(slide, [
    "Transaksi: Payment, Refund, Commission, Payout",
    "Gateway: Midtrans & Xendit (polymorphic payable)",
    "Auto-generated: TRX-{10 karakter}",
    "Status: Pending, Success, Failed, Expired, Refunded",
    "Payout vendor: Pending → Processing → Paid/Rejected",
    "Rekening bank vendor untuk payout",
    "Commission settings: rate per kategori atau global",
    "Mark Paid action dengan timestamp & processor",
    "Fee transaksi otomatis terhitung",
], Inches(9.5), Inches(2.0), Inches(3.3), Pt(11), DARK_GRAY, Pt(5))


# ==================== SLIDE 8: PENGATURAN ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Mockup: Pengaturan & Multi Bahasa", "Konfigurasi sistem, konten, dan dukungan 8 bahasa")
add_bottom_bar(slide)
add_slide_number(slide, 10)

add_screenshot(slide, "settings.png", Inches(0.5), Inches(1.45), Inches(8.5), Inches(5.8))

explanation_box = add_rounded_rect(slide, Inches(9.3), Inches(1.45), Inches(3.7), Inches(5.8), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5))
add_rect(slide, Inches(9.3), Inches(1.45), Inches(3.7), Inches(0.4), fill_color=GRAY)
add_text(slide, "  Fitur Pengaturan", Inches(9.5), Inches(1.48), Inches(3.3), Inches(0.35),
         font_size=Pt(12), color=WHITE, bold=True)
add_multi_text(slide, [
    "Key-Value settings: group, key, value, type",
    "5 grup: general, commission, payment, email, social",
    "Tipe: string, number, boolean, json, file",
    "Cached untuk performa (cache invalidation on update)",
    "Multi Bahasa (8 bahasa):",
    "  EN, ID, ES, FR, JA, KO, ZH, AR",
    "Dukungan RTL untuk bahasa Arab",
    "Language switcher di admin panel",
    "Manajemen konten: Blog, Banner, FAQ, Testimonial, Gallery",
], Inches(9.5), Inches(2.0), Inches(3.3), Pt(11), DARK_GRAY, Pt(5))


# ==================== SLIDE 9: ERD ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Design Database (ERD)", "37 tabel dengan relasi polymorphic, soft deletes, dan Spatie MediaLibrary")
add_bottom_bar(slide)
add_slide_number(slide, 11)

add_screenshot(slide, "erd.png", Inches(0.5), Inches(1.45), Inches(12.3), Inches(5.8))


# ==================== SLIDE 10: TABEL UTAMA ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Penjelasan Tabel Utama", "Tabel-tabel inti yang menjadi fondasi sistem BrightDor")
add_bottom_bar(slide)
add_slide_number(slide, 12)

tables_info = [
    ("users", "Tabel pusat pengguna (admin, vendor, couple). Menyimpan kredensial, profil, tipe user, dan status akun.", "id, name, email, phone, user_type, status", BLACK),
    ("vendors", "Data vendor pernikahan: profil usaha, lokasi, kontak, rating, status approval, data bank.", "id, user_id, vendor_category_id, business_name, status, is_verified", CREAM),
    ("services", "Layanan yang ditawarkan vendor: nama, harga, diskon, kapasitas, fitur, status publikasi.", "id, vendor_id, name, price, discount_price, status", CREAM),
    ("bookings", "Pesanan layanan: kode booking, data acara, komponen keuangan, status, catatan.", "id, booking_code, user_id, vendor_id, event_date, total_amount, status", INFO),
    ("invitation_templates", "Template undangan digital: nama, harga, preview, fitur, status premium.", "id, name, price, is_premium, is_active, sales_count", PURPLE),
    ("transactions", "Semua transaksi keuangan: payment, refund, commission, payout dengan polymorphic relation.", "id, transaction_code, type, amount, fee, net_amount, status", SUCCESS),
]
y = Inches(1.5)
for tname, tdesc, tcols, tcolor in tables_info:
    tcard = add_rounded_rect(slide, Inches(0.7), y, Inches(12.0), Inches(0.85), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5))
    add_rect(slide, Inches(0.7), y, Inches(0.08), Inches(0.85), fill_color=tcolor)
    add_text(slide, tname, Inches(1.0), y + Inches(0.05), Inches(2.2), Inches(0.3),
             font_size=Pt(13), color=BLACK, bold=True)
    add_text(slide, tdesc, Inches(3.3), y + Inches(0.05), Inches(6.5), Inches(0.35),
             font_size=Pt(10), color=GRAY)
    add_text(slide, tcols, Inches(3.3), y + Inches(0.45), Inches(9), Inches(0.3),
             font_size=Pt(9), color=LIGHT_GRAY)
    y += Inches(0.92)


# ==================== SLIDE 11: RELASI ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Relasi Antar Tabel", "Hubungan antar entitas dalam database BrightDor")
add_bottom_bar(slide)
add_slide_number(slide, 13)

relations = [
    ("User → Vendor", "HasOne", "1 user memiliki 1 vendor"),
    ("User → Booking", "HasMany", "1 user memiliki banyak booking"),
    ("User → InvitationOrder", "HasMany", "1 user memiliki banyak order undangan"),
    ("User → Transaction", "HasMany", "1 user memiliki banyak transaksi"),
    ("Vendor → Service", "HasMany", "1 vendor memiliki banyak layanan"),
    ("Vendor → Booking", "HasMany", "1 vendor menerima banyak booking"),
    ("Vendor → Payout", "HasMany", "1 vendor menerima banyak payout"),
    ("VendorCategory → Vendor", "HasMany", "1 kategori memiliki banyak vendor"),
    ("Booking → Transaction", "MorphMany", "1 booking bisa punya banyak transaksi"),
    ("InvitationOrder → Invitation", "HasOne", "1 order menghasilkan 1 invitation aktif"),
    ("InvitationOrder → Transaction", "MorphMany", "1 order undangan punya banyak transaksi"),
    ("Invitation → InvitationRsvp", "HasMany", "1 undangan menerima banyak RSVP"),
]
y = Inches(1.5)
for i, (rel, rtype, rdesc) in enumerate(relations):
    col = i % 2
    rx = Inches(0.7) + Inches(col * 6.3)
    ry = y + Inches((i // 2) * 0.72)
    rbox = add_rounded_rect(slide, rx, ry, Inches(5.9), Inches(0.6), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5))
    add_text(slide, rel, rx + Inches(0.15), ry + Inches(0.05), Inches(2.5), Inches(0.25),
             font_size=Pt(12), color=BLACK, bold=True)
    badge = add_rounded_rect(slide, rx + Inches(3.0), ry + Inches(0.08), Inches(1.2), Inches(0.22), fill_color=CREAM)
    add_text(slide, rtype, rx + Inches(3.0), ry + Inches(0.08), Inches(1.2), Inches(0.22),
             font_size=Pt(8), color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, rdesc, rx + Inches(0.15), ry + Inches(0.32), Inches(5.6), Inches(0.25),
             font_size=Pt(10), color=GRAY)


# ==================== SLIDE 12: ALUR BISNIS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Alur Bisnis Singkat", "Alur kerja utama dari marketplace BrightDor")
add_bottom_bar(slide)
add_slide_number(slide, 14)

flows = [
    ("1", "Vendor\nMendaftar", "Vendor buat akun,\nisi profil usaha", CREAM),
    ("2", "Admin\nReview", "Admin approve/reject\nvendor", INFO),
    ("3", "Vendor Upload\nLayanan", "Buat service\n& harga", PURPLE),
    ("4", "Couple\nBooking", "Pilih vendor\n& layanan", SUCCESS),
    ("5", "Pembayaran", "Bayar → transaksi\n& komisi otomatis", WARNING),
    ("6", "Payout\nVendor", "Admin proses\npayout vendor", RGBColor(0xEA, 0x58, 0x0C)),
]
x = Inches(0.7)
for num, title, desc, color in flows:
    fbox = add_rounded_rect(slide, x, Inches(1.7), Inches(1.9), Inches(2.6), fill_color=WHITE, border_color=color)
    circle = add_rounded_rect(slide, x + Inches(0.7), Inches(1.9), Inches(0.5), Inches(0.5), fill_color=color)
    add_text(slide, num, x + Inches(0.7), Inches(1.92), Inches(0.5), Inches(0.5),
             font_size=Pt(18), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.1), Inches(2.5), Inches(1.7), Inches(0.6),
             font_size=Pt(11), color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, desc, x + Inches(0.1), Inches(3.15), Inches(1.7), Inches(0.9),
             font_size=Pt(9), color=GRAY, alignment=PP_ALIGN.CENTER)
    if num != "6":
        add_text(slide, "→", x + Inches(1.95), Inches(2.6), Inches(0.3), Inches(0.4),
                 font_size=Pt(18), color=CREAM, bold=True, alignment=PP_ALIGN.CENTER)
    x += Inches(2.1)

add_text(slide, "Alur Marketplace", Inches(0.7), Inches(1.45), Inches(3), Inches(0.25),
         font_size=Pt(11), color=BLACK, bold=True)

inv_flows = [
    ("1", "Pilih\nTemplate", "Pilih template\nundangan", PURPLE),
    ("2", "Bayar\nTemplate", "Bayar → status\nPaid", WARNING),
    ("3", "Isi\nData", "Nama, tanggal,\nvenue, foto", INFO),
    ("4", "Publish", "Live di\nsubdomain", SUCCESS),
    ("5", "RSVP", "Tamu RSVP\nhadir/tidak", RGBColor(0xEA, 0x58, 0x0C)),
]
x = Inches(1.2)
for num, title, desc, color in inv_flows:
    fbox = add_rounded_rect(slide, x, Inches(4.7), Inches(2.1), Inches(2.3), fill_color=WHITE, border_color=color)
    circle = add_rounded_rect(slide, x + Inches(0.8), Inches(4.85), Inches(0.45), Inches(0.45), fill_color=color)
    add_text(slide, num, x + Inches(0.8), Inches(4.87), Inches(0.45), Inches(0.45),
             font_size=Pt(16), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.1), Inches(5.4), Inches(1.9), Inches(0.5),
             font_size=Pt(11), color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, desc, x + Inches(0.1), Inches(5.9), Inches(1.9), Inches(0.8),
             font_size=Pt(9), color=GRAY, alignment=PP_ALIGN.CENTER)
    if num != "5":
        add_text(slide, "→", x + Inches(2.15), Inches(5.4), Inches(0.3), Inches(0.3),
                 font_size=Pt(16), color=CREAM, bold=True, alignment=PP_ALIGN.CENTER)
    x += Inches(2.3)

add_text(slide, "Alur Undangan Digital", Inches(1.0), Inches(4.45), Inches(3), Inches(0.25),
         font_size=Pt(11), color=BLACK, bold=True)


# ==================== SLIDE 13: KESIMPULAN ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Kesimpulan")
add_bottom_bar(slide)
add_slide_number(slide, 15)

summary_items = [
    ("Laravel 11 + Filament v5", "Framework modern dengan admin panel powerful dan customisable"),
    ("37 Tabel Database", "Desain relasi komprehensif dengan polymorphic relations dan soft deletes"),
    ("7 Modul Utama", "Dashboard, Vendor, Marketplace, Undangan Digital, Keuangan, Konten, Pengaturan"),
    ("Multi Bahasa", "8 bahasa termasuk Indonesia dengan dukungan RTL untuk Arab"),
    ("Top Navigation Layout", "Desain clean tanpa sidebar, warna putih + cream + hitam"),
    ("Ekosistem Lengkap", "Dari pendaftaran vendor hingga payout, semuanya terintegrasi"),
]
y = Inches(1.5)
for title, desc in summary_items:
    sbox = add_rounded_rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.75), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5))
    add_rect(slide, Inches(0.8), y, Inches(0.06), Inches(0.75), fill_color=CREAM)
    add_text(slide, title, Inches(1.1), y + Inches(0.08), Inches(4), Inches(0.3),
             font_size=Pt(14), color=BLACK, bold=True)
    add_text(slide, desc, Inches(1.1), y + Inches(0.4), Inches(11.2), Inches(0.3),
             font_size=Pt(11), color=GRAY)
    y += Inches(0.85)

tech_box = add_rounded_rect(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5), fill_color=CREAM)
add_text(slide, "  Tech Stack: Laravel 11.31  |  Filament v5.0  |  Spatie Permission  |  Spatie MediaLibrary  |  Vite  |  Tailwind CSS  |  SQLite",
         Inches(0.9), Inches(6.72), Inches(11.5), Inches(0.4),
         font_size=Pt(11), color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)


# ==================== SLIDE 14: TERIMA KASIH ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_rect(slide, Inches(0), Inches(0), Inches(0.12), prs.slide_height, fill_color=CREAM)
add_rect(slide, Inches(0), Inches(3.4), prs.slide_width, Inches(0.005), fill_color=CREAM)

add_text(slide, "Terima Kasih", Inches(1.2), Inches(2.0), Inches(11), Inches(0.9),
         font_size=Pt(48), color=CREAM, bold=True)
add_text(slide, "Questions & Answers", Inches(1.2), Inches(2.8), Inches(11), Inches(0.5),
         font_size=Pt(20), color=WHITE)

add_line(slide, Inches(1.2), Inches(3.6), Inches(3), CREAM)

info_lines = [
    "BrightDor — Marketplace Jasa Pernikahan Premium",
    "Laravel 11 + Filament v5  |  37 Tabel  |  8 Bahasa",
    "Admin Panel: /admin  |  admin@brightdor.test",
]
y = Inches(4.0)
for line in info_lines:
    add_text(slide, line, Inches(1.2), y, Inches(11), Inches(0.35),
             font_size=Pt(14), color=LIGHT_GRAY)
    y += Inches(0.35)

add_rect(slide, Inches(11.5), Inches(0), Inches(1.833), prs.slide_height, fill_color=CREAM)
add_text(slide, "Q&A", Inches(11.6), Inches(3.2), Inches(1.6), Inches(0.6),
         font_size=Pt(28), color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, "2026", Inches(11.6), Inches(3.8), Inches(1.6), Inches(0.4),
         font_size=Pt(14), color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# Save
output_path = OUTPUT_PPT
prs.save(output_path)
print(f"PPT saved to: {output_path}")
# Also save to docs for submission
try:
    prs.save(OUTPUT_PPT_DOCS)
    print(f"PPT also saved to: {OUTPUT_PPT_DOCS}")
except Exception as e:
    print(f"Docs save failed: {e}")
