from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

CREAM = RGBColor(0xC4, 0xA5, 0x74)
CREAM_LIGHT = RGBColor(0xF5, 0xF0, 0xE8)
CREAM_BG = RGBColor(0xFA, 0xF7, 0xF2)
BLACK = RGBColor(0x14, 0x14, 0x14)
DARK_GRAY = RGBColor(0x2A, 0x2A, 0x2A)
GRAY = RGBColor(0x52, 0x52, 0x5B)
LIGHT_GRAY = RGBColor(0xA1, 0xA1, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x16, 0xA3, 0x4A)
WARNING = RGBColor(0xEA, 0xB3, 0x08)
DANGER = RGBColor(0xDC, 0x26, 0x26)
INFO = RGBColor(0x25, 0x63, 0xEB)

FONT_MAIN = "Plus Jakarta Sans"
FONT_FALLBACK = "Calibri"

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0), corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or WHITE
    line = shape.line
    if border_color:
        line.color.rgb = border_color
        line.width = border_width
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or WHITE
    line = shape.line
    if border_color:
        line.color.rgb = border_color
        line.width = border_width
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, text, left, top, width, height, font_size=Pt(14), color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_MAIN):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, items, left, top, width, font_size=Pt(13), color=BLACK, spacing=Pt(6), bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        prefix = "  •  " if bullet else ""
        p.text = prefix + item
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = FONT_MAIN
        p.space_after = spacing
        p.space_before = Pt(2)
    return txBox

def add_line(slide, left, top, width, color=CREAM, line_width=Pt(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_top_bar(slide):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=CREAM)

def add_bottom_bar(slide):
    add_rect(slide, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), fill_color=CREAM)

def add_slide_number(slide, num, total=15):
    add_text(slide, f"{num} / {total}", Inches(12.2), Inches(7.05), Inches(1), Inches(0.35),
             font_size=Pt(10), color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, title, subtitle=""):
    add_top_bar(slide)
    add_text(slide, title, Inches(0.8), Inches(0.35), Inches(10), Inches(0.55),
             font_size=Pt(28), color=BLACK, bold=True)
    add_line(slide, Inches(0.8), Inches(0.95), Inches(1.5), CREAM)
    if subtitle:
        add_text(slide, subtitle, Inches(0.8), Inches(1.05), Inches(10), Inches(0.4),
                 font_size=Pt(13), color=GRAY)

def add_mockup_box(slide, left, top, width, height, title, items, icon_text=""):
    box = add_shape(slide, left, top, width, height, fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
    add_rect(slide, left, top, width, Inches(0.45), fill_color=BLACK)
    add_text(slide, f"  {icon_text}  {title}" if icon_text else f"  {title}",
             left + Inches(0.1), top + Inches(0.05), width - Inches(0.2), Inches(0.4),
             font_size=Pt(12), color=WHITE, bold=True)
    y_offset = top + Inches(0.55)
    for item in items:
        add_text(slide, f"  {item}", left + Inches(0.2), y_offset, width - Inches(0.4), Inches(0.3),
                 font_size=Pt(10), color=DARK_GRAY)
        y_offset += Inches(0.28)
    return box

def add_stat_card(slide, left, top, label, value, color=CREAM):
    card = add_shape(slide, left, top, Inches(2.2), Inches(1.15), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
    add_rect(slide, left, top + Inches(0.95), Inches(2.2), Inches(0.06), fill_color=color)
    add_text(slide, value, left + Inches(0.15), top + Inches(0.12), Inches(1.9), Inches(0.45),
             font_size=Pt(22), color=BLACK, bold=True)
    add_text(slide, label, left + Inches(0.15), top + Inches(0.55), Inches(1.9), Inches(0.35),
             font_size=Pt(11), color=GRAY)

def add_nav_mockup(slide, left, top, width):
    add_rect(slide, left, top, width, Inches(0.6), fill_color=BLACK)
    add_text(slide, "  BrightDor", left + Inches(0.2), top + Inches(0.1), Inches(2), Inches(0.4),
             font_size=Pt(14), color=CREAM, bold=True)
    nav_items = ["Dashboard", "Vendors", "Marketplace", "Invitations", "Finance", "Content", "Settings"]
    x = left + width - Inches(4.5)
    for item in nav_items:
        clr = CREAM if item == "Dashboard" else WHITE
        add_text(slide, item, x, top + Inches(0.12), Inches(1.3), Inches(0.35),
                 font_size=Pt(10), color=clr, alignment=PP_ALIGN.CENTER)
        x += Inches(1.25)


# ==================== SLIDE 1: COVER ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_rect(slide, Inches(0), Inches(0), Inches(0.12), prs.slide_height, fill_color=CREAM)
add_rect(slide, Inches(0), Inches(3.2), prs.slide_width, Inches(0.005), fill_color=CREAM)

add_text(slide, "BRIGHTDOR", Inches(1.2), Inches(1.8), Inches(11), Inches(1),
         font_size=Pt(52), color=CREAM, bold=True, alignment=PP_ALIGN.LEFT)
add_text(slide, "Documentation Project", Inches(1.2), Inches(2.55), Inches(11), Inches(0.6),
         font_size=Pt(22), color=WHITE, alignment=PP_ALIGN.LEFT)

add_text(slide, "Marketplace Jasa Pernikahan Premium & Undangan Digital", Inches(1.2), Inches(3.5), Inches(11), Inches(0.5),
         font_size=Pt(16), color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

info_items = [
    "Laravel 11  |  Filament v5  |  SQLite",
    "Teknologi Informasi - Institut Teknologi"
]
y = Inches(5.2)
for item in info_items:
    add_text(slide, item, Inches(1.2), y, Inches(11), Inches(0.35),
             font_size=Pt(13), color=LIGHT_GRAY)
    y += Inches(0.35)

add_rect(slide, Inches(11.5), Inches(0), Inches(1.833), prs.slide_height, fill_color=CREAM)
add_text(slide, "2026", Inches(11.6), Inches(3.2), Inches(1.6), Inches(0.6),
         font_size=Pt(28), color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)


# ==================== SLIDE 2: NAMA APLIKASI ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "BrightDor", "Marketplace Jasa Pernikahan Premium & Undangan Digital")
add_bottom_bar(slide)
add_slide_number(slide, 2)

card = add_shape(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.2), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_text(slide, "Tentang BrightDor", Inches(1.1), Inches(1.75), Inches(5), Inches(0.4),
         font_size=Pt(18), color=BLACK, bold=True)
desc = (
    "BrightDor adalah platform digital yang menghubungkan pasangan pengantin (couple) "
    "dengan vendor jasa pernikahan premium di Indonesia. Platform ini menyediakan "
    "marketplace lengkap untuk layanan pernikahan beserta fitur undangan digital "
    "interaktif dengan sistem RSVP online."
)
add_text(slide, desc, Inches(1.1), Inches(2.2), Inches(11.2), Inches(1.2),
         font_size=Pt(13), color=GRAY)

features = [
    ("Vendor Management", "Approval, kategori, & verifikasi vendor"),
    ("Marketplace Jasa", "Katalog layanan pernikahan premium"),
    ("Booking & Order", "Sistem pemesanan & tracking status"),
    ("Undangan Digital", "Template premium + RSVP online"),
    ("Sistem Keuangan", "Transaksi, komisi, & payout vendor"),
    ("Multi Bahasa", "8 bahasa termasuk Indonesia & Arab"),
]
x = Inches(0.8)
y_start = Inches(4.1)
for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    fx = Inches(0.8) + Inches(col * 4.0)
    fy = y_start + Inches(row * 1.5)
    fcard = add_shape(slide, fx, fy, Inches(3.7), Inches(1.2), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
    add_rect(slide, fx, fy, Inches(3.7), Inches(0.05), fill_color=CREAM)
    add_text(slide, title, fx + Inches(0.15), fy + Inches(0.15), Inches(3.4), Inches(0.35),
             font_size=Pt(13), color=BLACK, bold=True)
    add_text(slide, desc, fx + Inches(0.15), fy + Inches(0.55), Inches(3.4), Inches(0.55),
             font_size=Pt(11), color=GRAY)


# ==================== SLIDE 3: TUJUAN & MANFAAT ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Tujuan & Manfaat Aplikasi")
add_bottom_bar(slide)
add_slide_number(slide, 3)

add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_rect(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.5), fill_color=BLACK)
add_text(slide, "  Tujuan Pengembangan", Inches(1.0), Inches(1.55), Inches(5.4), Inches(0.4),
         font_size=Pt(14), color=CREAM, bold=True)
tujuan = [
    "Menyediakan platform marketplace terpusat untuk jasa pernikahan premium",
    "Memudahkan pasangan pengantin menemukan vendor terpercaya",
    "Menyediakan solusi undangan digital modern dan interaktif",
    "Membantu vendor pernikahan memperluas jangkauan pasar",
    "Mengelola transaksi keuangan secara transparan dan efisien",
]
add_multi_text(slide, tujuan, Inches(1.1), Inches(2.2), Inches(5.3), Pt(12), DARK_GRAY, Pt(10), bullet=True)

add_shape(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.2), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_rect(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.5), fill_color=CREAM)
add_text(slide, "  Manfaat Aplikasi", Inches(7.1), Inches(1.55), Inches(5.4), Inches(0.4),
         font_size=Pt(14), color=BLACK, bold=True)
manfaat = [
    "Couple: Cari vendor, booking, buat undangan digital dalam satu platform",
    "Vendor: Kelola bisnis, terima booking, terima pembayaran",
    "Admin: Kontrol penuh atas seluruh ekosistem marketplace",
    "Keuangan: Sistem komisi otomatis & payout terintegrasi",
    "Global: Multi bahasa dengan dukungan RTL (Arab)",
]
add_multi_text(slide, manfaat, Inches(7.2), Inches(2.2), Inches(5.3), Pt(12), DARK_GRAY, Pt(10), bullet=True)


# ==================== SLIDE 4: MOCKUP DASHBOARD ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Mockup: Dashboard Admin", "Panel statistik utama untuk monitoring seluruh aktivitas marketplace")
add_bottom_bar(slide)
add_slide_number(slide, 4)

add_nav_mockup(slide, Inches(0.5), Inches(1.5), Inches(12.3))

add_stat_card(slide, Inches(0.7), Inches(2.3), "Total Vendors", "156", CREAM)
add_stat_card(slide, Inches(3.1), Inches(2.3), "Total Bookings", "1,284", INFO)
add_stat_card(slide, Inches(5.5), Inches(2.3), "Revenue", "Rp 4.2M", SUCCESS)
add_stat_card(slide, Inches(7.9), Inches(2.3), "Couples", "3,891", RGBColor(0x8B, 0x5C, 0xF6))
add_stat_card(slide, Inches(10.3), Inches(2.3), "Undangan Digital", "892", WARNING)

chart_box = add_shape(slide, Inches(0.7), Inches(3.7), Inches(7.8), Inches(3.1), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_text(slide, "Revenue Chart (6 Bulan)", Inches(1.0), Inches(3.85), Inches(3), Inches(0.3),
         font_size=Pt(12), color=BLACK, bold=True)
add_text(slide, "Line chart menampilkan trend revenue dari transaksi successful\nselama 6 bulan terakhir", Inches(1.0), Inches(4.25), Inches(7.2), Inches(0.7),
         font_size=Pt(10), color=GRAY)
for i in range(6):
    bx = Inches(1.5) + Inches(i * 1.1)
    h = Inches(0.8 + (i * 0.3 if i < 4 else 0.3))
    add_rect(slide, bx, Inches(6.3) - h, Inches(0.5), h, fill_color=CREAM)
add_line(slide, Inches(1.0), Inches(6.35), Inches(7.2), LIGHT_GRAY)

cat_box = add_shape(slide, Inches(8.8), Inches(3.7), Inches(3.9), Inches(3.1), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_text(slide, "Bookings by Kategori", Inches(9.1), Inches(3.85), Inches(3), Inches(0.3),
         font_size=Pt(12), color=BLACK, bold=True)
categories = [("Venue", "35%"), ("Catering", "25%"), ("Dekorasi", "20%"), ("Foto/Video", "12%"), ("Lainnya", "8%")]
y = Inches(4.3)
colors_list = [CREAM, RGBColor(0xD4, 0xB8, 0x96), RGBColor(0xBE, 0x9E, 0x7A), GRAY, LIGHT_GRAY]
for i, (cat, pct) in enumerate(categories):
    add_text(slide, f"{cat}", Inches(9.1), y, Inches(1.8), Inches(0.25), font_size=Pt(10), color=DARK_GRAY)
    bar_w = float(pct.strip('%')) / 100 * Inches(1.8)
    add_rect(slide, Inches(10.9), y + Inches(0.05), int(bar_w), Inches(0.15), fill_color=colors_list[i])
    add_text(slide, pct, Inches(10.9) + int(bar_w) + Inches(0.1), y, Inches(0.5), Inches(0.25), font_size=Pt(10), color=GRAY)
    y += Inches(0.45)

explanation = add_shape(slide, Inches(0.7), Inches(6.95), Inches(12.0), Inches(0.35), fill_color=CREAM_LIGHT)
add_text(slide, "  Dashboard menampilkan 5 kartu statistik utama, chart revenue 6 bulan, distribusi booking per kategori, quick actions, dan aktivitas terbaru.",
         Inches(0.9), Inches(6.97), Inches(11.5), Inches(0.3), font_size=Pt(10), color=DARK_GRAY)


# ==================== SLIDE 5: MOCKUP MANAJEMEN VENDOR ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Mockup: Manajemen Vendor", "Approval, kategori, verifikasi, dan pengelolaan vendor pernikahan")
add_bottom_bar(slide)
add_slide_number(slide, 5)

add_nav_mockup(slide, Inches(0.5), Inches(1.5), Inches(12.3))

table_box = add_shape(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(3.0), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.4), fill_color=BLACK)
cols = [("Business Name", 2.5), ("Owner", 1.5), ("Kategori", 1.5), ("Kota", 1.2), ("Rating", 1.0), ("Status", 1.2), ("Verified", 0.8), ("Aksi", 1.5)]
x = Inches(0.9)
for col_name, w in cols:
    add_text(slide, col_name, x, Inches(2.33), Inches(w), Inches(0.35),
             font_size=Pt(10), color=CREAM, bold=True)
    x += Inches(w + 0.1)

vendors_data = [
    ("Elegant Venue Bandung", "Rina Sari", "Venue", "Bandung", "4.8", "Approved", True),
    ("Catering Masak Jawa", "Budi Santoso", "Catering", "Jakarta", "4.5", "Pending", False),
    ("Dekorasi Bunga Indah", "Siti Nurhaliza", "Dekorasi", "Surabaya", "4.9", "Approved", True),
    ("Photo Studio Pro", "Andi Wijaya", "Foto/Video", "Yogyakarta", "4.3", "Rejected", False),
]
y = Inches(2.85)
for vname, owner, cat, city, rating, status, verified in vendors_data:
    x = Inches(0.9)
    vals = [vname, owner, cat, city, f"★ {rating}", status, "✓" if verified else "—"]
    for i, (val, (_, w)) in enumerate(zip(vals, cols)):
        clr = DARK_GRAY
        if i == 5:
            clr = SUCCESS if status == "Approved" else (WARNING if status == "Pending" else DANGER)
        if i == 6:
            clr = SUCCESS if verified else LIGHT_GRAY
        bld = i == 0
        add_text(slide, str(val), x, y, Inches(w), Inches(0.28),
                 font_size=Pt(9), color=clr, bold=bld)
        x += Inches(w + 0.1)
    y += Inches(0.35)

actions_box = add_shape(slide, Inches(0.7), Inches(5.5), Inches(5.8), Inches(1.55), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_text(slide, "Aksi yang Tersedia:", Inches(0.9), Inches(5.6), Inches(3), Inches(0.3),
         font_size=Pt(12), color=BLACK, bold=True)
add_multi_text(slide, [
    "Approve: Set status approved, tandai verified",
    "Reject: Input alasan penolakan",
    "Suspend: Nonaktifkan vendor aktif",
    "View: Lihat detail lengkap vendor",
], Inches(0.9), Inches(5.95), Inches(5.4), Pt(11), DARK_GRAY, Pt(5), bullet=True)

cat_box = add_shape(slide, Inches(6.8), Inches(5.5), Inches(5.9), Inches(1.55), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_text(slide, "Kategori Vendor:", Inches(7.0), Inches(5.6), Inches(3), Inches(0.3),
         font_size=Pt(12), color=BLACK, bold=True)
add_multi_text(slide, [
    "Venue, Catering, Dekorasi, Foto/Video",
    "Musik & Entertainment, Makeup & Busana",
    "Wedding Organizer, Lainnya",
    "Commission rate per kategori dapat diatur",
], Inches(7.0), Inches(5.95), Inches(5.5), Pt(11), DARK_GRAY, Pt(5), bullet=True)


# ==================== SLIDE 6: MOCKUP BOOKING/ORDER ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Mockup: Booking & Order", "Sistem pemesanan layanan vendor dengan tracking status real-time")
add_bottom_bar(slide)
add_slide_number(slide, 6)

add_nav_mockup(slide, Inches(0.5), Inches(1.5), Inches(12.3))

table_box = add_shape(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(2.8), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.4), fill_color=BLACK)
bcols = [("Kode Booking", 1.6), ("Customer", 1.5), ("Vendor", 1.8), ("Layanan", 1.8), ("Tanggal", 1.2), ("Total", 1.3), ("Status", 1.2)]
x = Inches(0.9)
for col_name, w in bcols:
    add_text(slide, col_name, x, Inches(2.33), Inches(w), Inches(0.35),
             font_size=Pt(10), color=CREAM, bold=True)
    x += Inches(w + 0.1)

bookings_data = [
    ("BD-A1B2C3D4", "Rina & Andi", "Elegant Venue", "Paket Premium", "25 Agt 2026", "Rp 15.000.000", "Confirmed"),
    ("BD-E5F6G7H8", "Sinta & Budi", "Catering Masak", "Paket Lengkap", "12 Sep 2026", "Rp 8.500.000", "Pending"),
    ("BD-I9J0K1L2", "Maya & Raka", "Dekorasi Bunga", "Paket Standard", "05 Okt 2026", "Rp 5.000.000", "On Progress"),
    ("BD-M3N4O5P6", "Dewi & Fajar", "Photo Studio", "Paket Eksklusif", "18 Jul 2026", "Rp 12.000.000", "Completed"),
]
y = Inches(2.85)
for bdata in bookings_data:
    x = Inches(0.9)
    for i, (val, (_, w)) in enumerate(zip(bdata, bcols)):
        clr = DARK_GRAY
        if i == 6:
            clr = SUCCESS if val == "Completed" else (WARNING if val == "Pending" else (INFO if val == "Confirmed" else RGBColor(0x8B, 0x5C, 0xF6)))
        bld = i == 0
        add_text(slide, val, x, y, Inches(w), Inches(0.28),
                 font_size=Pt(9), color=clr, bold=bld)
        x += Inches(w + 0.1)
    y += Inches(0.35)

detail_box = add_shape(slide, Inches(0.7), Inches(5.3), Inches(5.8), Inches(1.8), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_text(slide, "Detail Booking:", Inches(0.9), Inches(5.4), Inches(3), Inches(0.3),
         font_size=Pt(12), color=BLACK, bold=True)
add_multi_text(slide, [
    "Status flow: Pending → Confirmed → On Progress → Completed",
    "Opsi: Cancel / Refund dengan alasan",
    "Tanggal & waktu acara, lokasi, jumlah tamu",
    "Catatan customer & catatan admin",
], Inches(0.9), Inches(5.75), Inches(5.4), Pt(11), DARK_GRAY, Pt(5), bullet=True)

finance_box = add_shape(slide, Inches(6.8), Inches(5.3), Inches(5.9), Inches(1.8), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_text(slide, "Komponen Keuangan:", Inches(7.0), Inches(5.4), Inches(3), Inches(0.3),
         font_size=Pt(12), color=BLACK, bold=True)
add_multi_text(slide, [
    "Subtotal, diskon, admin fee",
    "Komisi otomatis berdasarkan kategori vendor",
    "Total amount yang harus dibayar customer",
    "Tracking transaksi terkait (polymorphic)",
], Inches(7.0), Inches(5.75), Inches(5.5), Pt(11), DARK_GRAY, Pt(5), bullet=True)


# ==================== SLIDE 7: MOCKUP UNDANGAN DIGITAL ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Mockup: Undangan Digital", "Template premium undangan digital + RSVP online interaktif")
add_bottom_bar(slide)
add_slide_number(slide, 7)

add_nav_mockup(slide, Inches(0.5), Inches(1.5), Inches(12.3))

template_cards = [
    ("Elegant Gold", "Rp 250.000", "Premium", True),
    ("Garden Romance", "Rp 150.000", "Standard", False),
    ("Modern Minimalis", "Rp 200.000", "Premium", True),
    ("Classic Rose", "Rp 175.000", "Standard", False),
]
x = Inches(0.7)
for name, price, tipe, premium in template_cards:
    tcard = add_shape(slide, x, Inches(2.2), Inches(2.85), Inches(2.0), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
    add_rect(slide, x, Inches(2.2), Inches(2.85), Inches(1.0), fill_color=CREAM_LIGHT)
    add_text(slide, f"  {name}", x + Inches(0.1), Inches(2.6), Inches(2.6), Inches(0.4),
             font_size=Pt(14), color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    badge_color = WARNING if premium else LIGHT_GRAY
    badge = add_shape(slide, x + Inches(0.8), Inches(3.35), Inches(1.2), Inches(0.3), fill_color=badge_color)
    add_text(slide, tipe, x + Inches(0.8), Inches(3.35), Inches(1.2), Inches(0.3),
             font_size=Pt(9), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, price, x + Inches(0.1), Inches(3.75), Inches(2.6), Inches(0.3),
             font_size=Pt(12), color=CREAM, bold=True, alignment=PP_ALIGN.CENTER)
    x += Inches(3.1)

order_box = add_shape(slide, Inches(0.7), Inches(4.4), Inches(12.0), Inches(1.2), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_text(slide, "Order Undangan Digital", Inches(0.9), Inches(4.5), Inches(5), Inches(0.3),
         font_size=Pt(13), color=BLACK, bold=True)
add_multi_text(slide, [
    "Status: Pending → Paid → Active → Expired / Cancelled",
    "Data: Nama Pengantin, Tanggal Nikah, Venue, Subdomain, Custom Domain",
    "Invitation: Content JSON, Theme Settings, RSVP Counter (Yes/No/Maybe)",
], Inches(0.9), Inches(4.85), Inches(11.5), Pt(11), DARK_GRAY, Pt(5), bullet=True)

rsvp_box = add_shape(slide, Inches(0.7), Inches(5.8), Inches(5.8), Inches(1.3), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_text(slide, "RSVP System:", Inches(0.9), Inches(5.9), Inches(3), Inches(0.3),
         font_size=Pt(12), color=BLACK, bold=True)
add_multi_text(slide, [
    "Tamu dapat RSVP: Hadir / Tidak / Mungkin",
    "Input: Nama, Email, Telepon, Jumlah Tamu, Pesan",
    "Counter real-time pada invitation page",
], Inches(0.9), Inches(6.2), Inches(5.4), Pt(11), DARK_GRAY, Pt(5), bullet=True)

domain_box = add_shape(slide, Inches(6.8), Inches(5.8), Inches(5.9), Inches(1.3), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_text(slide, "Domain Management:", Inches(7.0), Inches(5.9), Inches(3), Inches(0.3),
         font_size=Pt(12), color=BLACK, bold=True)
add_multi_text(slide, [
    "Subdomain otomatis: {nama}.brightdor.id",
    "Custom domain support (unique constraint)",
    "Published/unpublished toggle",
], Inches(7.0), Inches(6.2), Inches(5.5), Pt(11), DARK_GRAY, Pt(5), bullet=True)


# ==================== SLIDE 8: MOCKUP KEUANGAN ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Mockup: Keuangan", "Transaksi, komisi vendor, payout, dan pengaturan komisi")
add_bottom_bar(slide)
add_slide_number(slide, 8)

add_nav_mockup(slide, Inches(0.5), Inches(1.5), Inches(12.3))

trans_box = add_shape(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(2.2), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_rect(slide, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.4), fill_color=BLACK)
add_text(slide, "  Transaksi Terbaru", Inches(0.9), Inches(2.33), Inches(5), Inches(0.35),
         font_size=Pt(12), color=CREAM, bold=True)
tcols = [("Kode", 1.8), ("User", 1.5), ("Tipe", 1.2), ("Amount", 1.5), ("Gateway", 1.2), ("Status", 1.2)]
x = Inches(0.9)
for col_name, w in tcols:
    add_text(slide, col_name, x, Inches(2.75), Inches(w), Inches(0.28),
             font_size=Pt(9), color=GRAY, bold=True)
    x += Inches(w + 0.1)
trans_data = [
    ("TRX-1A2B3C4D5E", "Rina & Andi", "Payment", "Rp 15.000.000", "Midtrans", "Success"),
    ("TRX-6F7G8H9I0J", "Sinta & Budi", "Payment", "Rp 8.500.000", "Xendit", "Pending"),
    ("TRX-K1L2M3N4O5", "Admin", "Payout", "Rp 3.200.000", "Manual", "Processing"),
]
y = Inches(3.05)
for tdata in trans_data:
    x = Inches(0.9)
    for i, (val, (_, w)) in enumerate(zip(tdata, tcols)):
        clr = DARK_GRAY
        if i == 5:
            clr = SUCCESS if val == "Success" else (WARNING if val == "Pending" else INFO)
        add_text(slide, val, x, y, Inches(w), Inches(0.25),
                 font_size=Pt(9), color=clr, bold=(i==0))
        x += Inches(w + 0.1)
    y += Inches(0.3)

payout_box = add_shape(slide, Inches(0.7), Inches(4.7), Inches(5.8), Inches(2.3), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_text(slide, "Payout Vendor", Inches(0.9), Inches(4.8), Inches(3), Inches(0.3),
         font_size=Pt(13), color=BLACK, bold=True)
add_multi_text(slide, [
    "Status: Pending → Processing → Paid / Rejected",
    "Data: Kode payout, vendor, jumlah, fee, net amount",
    "Rekening bank: nama bank, nomor, nama rekening",
    "Admin notes & processor tracking",
    "Mark Paid action dengan timestamp",
], Inches(0.9), Inches(5.2), Inches(5.4), Pt(11), DARK_GRAY, Pt(5), bullet=True)

comm_box = add_shape(slide, Inches(6.8), Inches(4.7), Inches(5.9), Inches(2.3), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_text(slide, "Pengaturan Komisi", Inches(7.0), Inches(4.8), Inches(3), Inches(0.3),
         font_size=Pt(13), color=BLACK, bold=True)
add_multi_text(slide, [
    "Rate persen per kategori vendor (default 10%)",
    "Rate flat (Rp) per kategori",
    "Global commission setting (nullable category)",
    "Toggle active/inactive per setting",
    "Label untuk identifikasi komisi",
], Inches(7.0), Inches(5.2), Inches(5.5), Pt(11), DARK_GRAY, Pt(5), bullet=True)


# ==================== SLIDE 9: MOCKUP PENGATURAN ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Mockup: Pengaturan & Multi Bahasa", "Konfigurasi sistem, konten, dan dukungan 8 bahasa")
add_bottom_bar(slide)
add_slide_number(slide, 9)

add_nav_mockup(slide, Inches(0.5), Inches(1.5), Inches(12.3))

setting_box = add_shape(slide, Inches(0.7), Inches(2.3), Inches(5.8), Inches(4.7), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_rect(slide, Inches(0.7), Inches(2.3), Inches(5.8), Inches(0.4), fill_color=BLACK)
add_text(slide, "  Pengaturan Sistem (Key-Value)", Inches(0.9), Inches(2.33), Inches(5.4), Inches(0.35),
         font_size=Pt(12), color=CREAM, bold=True)

groups = ["general", "commission", "payment", "email", "social"]
y = Inches(2.85)
for grp in groups:
    gbadge = add_shape(slide, Inches(0.9), y, Inches(1.5), Inches(0.3), fill_color=CREAM_LIGHT)
    add_text(slide, grp, Inches(0.9), y, Inches(1.5), Inches(0.3),
             font_size=Pt(10), color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, "key → value (type: string/number/boolean/json/file)", Inches(2.6), y, Inches(3.7), Inches(0.3),
             font_size=Pt(10), color=GRAY)
    y += Inches(0.4)

add_multi_text(slide, [
    "Key-value pairs untuk konfigurasi fleksibel",
    "Tipe data: string, number, boolean, json, file",
    "Cached untuk performa (cache invalidation on update)",
    "Grouping: general, commission, payment, email, social",
], Inches(0.9), Inches(4.9), Inches(5.4), Pt(11), DARK_GRAY, Pt(5), bullet=True)

content_box = add_shape(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(4.7), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_rect(slide, Inches(6.8), Inches(2.3), Inches(5.9), Inches(0.4), fill_color=CREAM)
add_text(slide, "  Manajemen Konten & Multi Bahasa", Inches(7.0), Inches(2.33), Inches(5.5), Inches(0.35),
         font_size=Pt(12), color=BLACK, bold=True)

content_items = [
    ("Blog", "Artikel dengan Rich Editor, status draft/published"),
    ("Banner", "Promosi home hero dengan scheduling"),
    ("Testimonial", "Review pasangan, rating 1-5, sort order"),
    ("FAQ", "Pertanyaan umum per kategori"),
    ("Gallery", "Galeri foto per vendor atau kategori"),
]
y = Inches(2.85)
for cname, cdesc in content_items:
    add_text(slide, f"▸ {cname}", Inches(7.0), y, Inches(2), Inches(0.28),
             font_size=Pt(11), color=BLACK, bold=True)
    add_text(slide, cdesc, Inches(8.5), y, Inches(4), Inches(0.28),
             font_size=Pt(10), color=GRAY)
    y += Inches(0.4)

add_line(slide, Inches(7.0), y + Inches(0.1), Inches(5.4), CREAM)

add_text(slide, "Multi Bahasa (8 Bahasa):", Inches(7.0), y + Inches(0.3), Inches(5.4), Inches(0.3),
         font_size=Pt(12), color=BLACK, bold=True)
langs = ["English", "Indonesia", "Español", "Français", "日本語", "한국어", "中文", "العربية"]
lang_y = y + Inches(0.65)
for i, lang in enumerate(langs):
    col = i % 4
    row = i // 4
    lx = Inches(7.0) + Inches(col * 1.4)
    ly = lang_y + Inches(row * 0.4)
    lbadge = add_shape(slide, lx, ly, Inches(1.2), Inches(0.3), fill_color=CREAM_LIGHT if i % 2 == 0 else WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
    add_text(slide, lang, lx, ly, Inches(1.2), Inches(0.3),
             font_size=Pt(10), color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# ==================== SLIDE 10: ERD OVERVIEW ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Overview Design Database (ERD)", "33 tabel dengan relasi polymorphic dan soft deletes")
add_bottom_bar(slide)
add_slide_number(slide, 10)

groups_erd = [
    ("Authentication", Inches(0.7), Inches(1.6), ["users", "sessions", "password_reset_tokens", "cache"], BLACK),
    ("Authorization (Spatie)", Inches(4.2), Inches(1.6), ["permissions", "roles", "model_has_permissions", "model_has_roles", "role_has_permissions"], GRAY),
    ("Vendor System", Inches(0.7), Inches(3.4), ["vendor_categories", "vendors", "services"], CREAM),
    ("Booking System", Inches(4.2), Inches(3.4), ["bookings"], INFO),
    ("Digital Invitations", Inches(6.5), Inches(1.6), ["invitation_template_categories", "invitation_templates", "invitation_orders", "invitations", "invitation_rsvps"], RGBColor(0x8B, 0x5C, 0xF6)),
    ("Finance", Inches(6.5), Inches(3.4), ["transactions", "commission_settings", "payouts"], SUCCESS),
    ("Content", Inches(9.5), Inches(1.6), ["blogs", "testimonials", "banners", "faqs", "galleries"], RGBColor(0xEA, 0x58, 0x0C)),
    ("System", Inches(9.5), Inches(3.4), ["settings", "activity_logs", "media", "jobs", "job_batches", "failed_jobs"], LIGHT_GRAY),
]
for gname, gx, gy, tables, gcolor in groups_erd:
    gbox = add_shape(slide, gx, gy, Inches(3.2), Inches(1.5) if len(tables) <= 4 else Inches(1.8),
                     fill_color=WHITE, border_color=gcolor, border_width=Pt(1.5))
    add_rect(slide, gx, gy, Inches(3.2), Inches(0.35), fill_color=gcolor)
    add_text(slide, f"  {gname}", gx + Inches(0.1), gy + Inches(0.02), Inches(3), Inches(0.3),
             font_size=Pt(10), color=WHITE if gcolor not in [LIGHT_GRAY, CREAM] else BLACK, bold=True)
    ty = gy + Inches(0.4)
    for t in tables:
        add_text(slide, f"  {t}", gx + Inches(0.15), ty, Inches(3), Inches(0.22),
                 font_size=Pt(8), color=DARK_GRAY)
        ty += Inches(0.22)

info_box = add_shape(slide, Inches(0.7), Inches(5.5), Inches(12.0), Inches(1.6), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
add_text(slide, "Karakteristik Database:", Inches(0.9), Inches(5.6), Inches(3), Inches(0.3),
         font_size=Pt(13), color=BLACK, bold=True)
add_multi_text(slide, [
    "33 tabel total  |  20 Eloquent Models  |  2 Polymorphic Relations (Transaction.payable, ActivityLog.subject)",
    "Soft Deletes pada tabel: vendors, services, bookings, invitation_templates, invitation_orders, invitations, blogs",
    "Spatie MediaLibrary pada: Vendor (logo, portfolio, documents), Service (cover, gallery), InvitationTemplate (preview, gallery)",
    "Auto-generated codes: BD-{8} (booking), INV-{8} (invitation order), TRX-{10} (transaction), PO-{8} (payout)",
], Inches(0.9), Inches(5.95), Inches(11.5), Pt(11), DARK_GRAY, Pt(5), bullet=True)


# ==================== SLIDE 11: TABEL UTAMA ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Penjelasan Tabel Utama", "Tabel-tabel inti yang menjadi fondasi sistem")
add_bottom_bar(slide)
add_slide_number(slide, 11)

tables_info = [
    ("users", "Tabel pusaut pengguna (admin, vendor, couple). Menyimpan kredensial, profil, tipe user, dan status akun.", "id, name, email, phone, avatar, user_type, status"),
    ("vendors", "Data vendor pernikahan: profil usaha, lokasi, kontak, rating, status approval, data bank.", "id, user_id, vendor_category_id, business_name, status, is_verified"),
    ("services", "Layanan yang ditawarkan vendor: nama, harga, diskon, kapasitas, fitur, status publikasi.", "id, vendor_id, name, price, discount_price, status, is_active"),
    ("bookings", "Pesanan layanan: kode booking, data acara, komponen keuangan, status, catatan.", "id, booking_code, user_id, vendor_id, event_date, total_amount, status"),
    ("invitation_templates", "Template undangan digital: nama, harga, preview, fitur, status premium.", "id, name, price, is_premium, is_active, sales_count"),
    ("transactions", "Semua transaksi keuangan: payment, refund, commission, payout dengan polymorphic relation.", "id, transaction_code, type, amount, fee, net_amount, status"),
]
y = Inches(1.5)
for tname, tdesc, tcols in tables_info:
    tcard = add_shape(slide, Inches(0.7), y, Inches(12.0), Inches(0.85), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
    add_rect(slide, Inches(0.7), y, Inches(0.08), Inches(0.85), fill_color=CREAM)
    add_text(slide, tname, Inches(1.0), y + Inches(0.05), Inches(2.2), Inches(0.3),
             font_size=Pt(13), color=BLACK, bold=True)
    add_text(slide, tdesc, Inches(3.3), y + Inches(0.05), Inches(6.5), Inches(0.35),
             font_size=Pt(10), color=GRAY)
    add_text(slide, tcols, Inches(3.3), y + Inches(0.45), Inches(9), Inches(0.3),
             font_size=Pt(9), color=LIGHT_GRAY)
    y += Inches(0.92)


# ==================== SLIDE 12: RELASI ANTAR TABEL ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Relasi Antar Tabel", "Hubungan antar entitas dalam database BrightDor")
add_bottom_bar(slide)
add_slide_number(slide, 12)

relations = [
    ("User → Vendor", "HasOne", "1 user memiliki 1 vendor (vendor_type)"),
    ("User → Booking", "HasMany", "1 user (couple) memiliki banyak booking"),
    ("User → InvitationOrder", "HasMany", "1 user memiliki banyak order undangan"),
    ("User → Transaction", "HasMany", "1 user memiliki banyak transaksi"),
    ("Vendor → Service", "HasMany", "1 vendor memiliki banyak layanan"),
    ("Vendor → Booking", "HasMany", "1 vendor menerima banyak booking"),
    ("Vendor → Payout", "HasMany", "1 vendor menerima banyak payout"),
    ("VendorCategory → Vendor", "HasMany", "1 kategori memiliki banyak vendor"),
    ("Booking → Transaction", "MorphMany", "1 booking bisa punya banyak transaksi"),
    ("InvitationOrder → Invitation", "HasOne", "1 order menghasilkan 1 invitation aktif"),
    ("InvitationOrder → Transaction", "MorphMany", "1 order undangan punya banyak transaksi"),
    ("Invitation → InvitationRsvp", "HasMany", "1 undangan menerima banyak RSVP"),
]
y = Inches(1.5)
for i, (rel, rtype, rdesc) in enumerate(relations):
    col = i % 2
    rx = Inches(0.7) + Inches(col * 6.3)
    ry = y + Inches((i // 2) * 0.72)
    rbox = add_shape(slide, rx, ry, Inches(5.9), Inches(0.6), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
    add_text(slide, rel, rx + Inches(0.15), ry + Inches(0.05), Inches(2.5), Inches(0.25),
             font_size=Pt(12), color=BLACK, bold=True)
    badge = add_shape(slide, rx + Inches(3.0), ry + Inches(0.08), Inches(1.2), Inches(0.22), fill_color=CREAM)
    add_text(slide, rtype, rx + Inches(3.0), ry + Inches(0.08), Inches(1.2), Inches(0.22),
             font_size=Pt(8), color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, rdesc, rx + Inches(0.15), ry + Inches(0.32), Inches(5.6), Inches(0.25),
             font_size=Pt(10), color=GRAY)


# ==================== SLIDE 13: ALUR BISNIS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Alur Bisnis Singkat", "Alur kerja utama dari marketplace BrightDor")
add_bottom_bar(slide)
add_slide_number(slide, 13)

flows = [
    ("1", "Vendor Mendaftar", "Vendor membuat akun, mengisi profil\nusaha, dan mengajukan diri", CREAM),
    ("2", "Admin Review", "Admin approve/reject vendor.\nJika approve → vendor verified", INFO),
    ("3", "Vendor Upload\nLayanan", "Vendor membuat service/layanan\n dengan harga & deskripsi", RGBColor(0x8B, 0x5C, 0xF6)),
    ("4", "Couple Booking", "Couple memilih vendor & layanan,\nmelakukan pemesanan (booking)", SUCCESS),
    ("5", "Pembayaran", "Couple bayar → transaksi tercatat.\nKomisi otomatis dihitung", WARNING),
    ("6", "Payout Vendor", "Admin proses payout ke vendor.\nNet amount = total - komisi", RGBColor(0xEA, 0x58, 0x0C)),
]
x = Inches(0.7)
for num, title, desc, color in flows:
    fbox = add_shape(slide, x, Inches(1.7), Inches(1.9), Inches(2.8), fill_color=WHITE, border_color=color, border_width=Pt(1.5))
    circle = add_shape(slide, x + Inches(0.7), Inches(1.9), Inches(0.5), Inches(0.5), fill_color=color)
    add_text(slide, num, x + Inches(0.7), Inches(1.92), Inches(0.5), Inches(0.5),
             font_size=Pt(18), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.1), Inches(2.55), Inches(1.7), Inches(0.6),
             font_size=Pt(11), color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, desc, x + Inches(0.1), Inches(3.2), Inches(1.7), Inches(1.0),
             font_size=Pt(9), color=GRAY, alignment=PP_ALIGN.CENTER)
    if num != "6":
        add_text(slide, "→", x + Inches(1.95), Inches(2.8), Inches(0.3), Inches(0.4),
                 font_size=Pt(18), color=CREAM, bold=True, alignment=PP_ALIGN.CENTER)
    x += Inches(2.1)

inv_flow = [
    ("1", "Pilih Template", "Couple memilih template\nundangan digital", RGBColor(0x8B, 0x5C, 0xF6)),
    ("2", "Bayar Template", "Pembayaran template.\nStatus: Pending → Paid", WARNING),
    ("3", "Isi Data", "Nama pengantin, tanggal,\nvenue, upload foto", INFO),
    ("4", "Publish", "Undangan live di\nsubdomain/custom domain", SUCCESS),
    ("5", "RSVP", "Tamu RSVP hadir/tidak.\nCounter real-time", RGBColor(0xEA, 0x58, 0x0C)),
]
x = Inches(1.2)
for num, title, desc, color in inv_flow:
    fbox = add_shape(slide, x, Inches(4.8), Inches(2.1), Inches(2.2), fill_color=WHITE, border_color=color, border_width=Pt(1))
    circle = add_shape(slide, x + Inches(0.8), Inches(4.95), Inches(0.45), Inches(0.45), fill_color=color)
    add_text(slide, num, x + Inches(0.8), Inches(4.97), Inches(0.45), Inches(0.45),
             font_size=Pt(16), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.1), Inches(5.5), Inches(1.9), Inches(0.4),
             font_size=Pt(11), color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, desc, x + Inches(0.1), Inches(5.9), Inches(1.9), Inches(0.8),
             font_size=Pt(9), color=GRAY, alignment=PP_ALIGN.CENTER)
    if num != "5":
        add_text(slide, "→", x + Inches(2.15), Inches(5.5), Inches(0.3), Inches(0.3),
                 font_size=Pt(16), color=CREAM, bold=True, alignment=PP_ALIGN.CENTER)
    x += Inches(2.3)

add_text(slide, "Alur Marketplace:", Inches(0.7), Inches(4.55), Inches(3), Inches(0.25),
         font_size=Pt(11), color=BLACK, bold=True)
add_text(slide, "Alur Undangan Digital:", Inches(1.0), Inches(4.55), Inches(3), Inches(0.25),
         font_size=Pt(11), color=BLACK, bold=True)


# ==================== SLIDE 14: KESIMPULAN ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM_BG)
add_section_header(slide, "Kesimpulan")
add_bottom_bar(slide)
add_slide_number(slide, 14)

summary_items = [
    ("Laravel 11 + Filament v5", "Framework modern dengan admin panel powerful dan customisable"),
    ("33 Tabel Database", "Desain relasi yang komprehensif dengan polymorphic relations dan soft deletes"),
    ("7 Modul Utama", "Dashboard, Vendor, Marketplace, Undangan Digital, Keuangan, Konten, Pengaturan"),
    ("Multi Bahasa", "8 bahasa termasuk Indonesia dengan dukungan RTL untuk Arab"),
    ("Top Navigation Layout", "Desain clean tanpa sidebar, warna putih + cream + hitam"),
    ("Ekosistem Lengkap", "Dari pendaftaran vendor hingga payout, semuanya terintegrasi"),
]

y = Inches(1.5)
for title, desc in summary_items:
    sbox = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.75), fill_color=WHITE, border_color=RGBColor(0xE5, 0xE5, 0xE5), border_width=Pt(1))
    add_rect(slide, Inches(0.8), y, Inches(0.06), Inches(0.75), fill_color=CREAM)
    add_text(slide, title, Inches(1.1), y + Inches(0.08), Inches(4), Inches(0.3),
             font_size=Pt(14), color=BLACK, bold=True)
    add_text(slide, desc, Inches(1.1), y + Inches(0.4), Inches(11.2), Inches(0.3),
             font_size=Pt(11), color=GRAY)
    y += Inches(0.85)

tech_box = add_shape(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5), fill_color=CREAM)
add_text(slide, "  Tech Stack: Laravel 11.31  |  Filament v5.0  |  Spatie Permission  |  Spatie MediaLibrary  |  Vite  |  Tailwind CSS  |  SQLite",
         Inches(0.9), Inches(6.72), Inches(11.5), Inches(0.4),
         font_size=Pt(11), color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)


# ==================== SLIDE 15: TERIMA KASIH ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)

add_rect(slide, Inches(0), Inches(0), Inches(0.12), prs.slide_height, fill_color=CREAM)
add_rect(slide, Inches(0), Inches(3.4), prs.slide_width, Inches(0.005), fill_color=CREAM)

add_text(slide, "Terima Kasih", Inches(1.2), Inches(2.0), Inches(11), Inches(0.9),
         font_size=Pt(48), color=CREAM, bold=True, alignment=PP_ALIGN.LEFT)
add_text(slide, "Questions & Answers", Inches(1.2), Inches(2.8), Inches(11), Inches(0.5),
         font_size=Pt(20), color=WHITE, alignment=PP_ALIGN.LEFT)

add_line(slide, Inches(1.2), Inches(3.6), Inches(3), CREAM)

info_lines = [
    "BrightDor - Marketplace Jasa Pernikahan Premium",
    "Laravel 11 + Filament v5  |  33 Tabel  |  8 Bahasa",
    "Admin Panel: /admin  |  admin@brightdor.test",
]
y = Inches(4.0)
for line in info_lines:
    add_text(slide, line, Inches(1.2), y, Inches(11), Inches(0.35),
             font_size=Pt(14), color=LIGHT_GRAY)
    y += Inches(0.35)

add_rect(slide, Inches(11.5), Inches(0), Inches(1.833), prs.slide_height, fill_color=CREAM)
add_text(slide, "Q&A", Inches(11.6), Inches(3.2), Inches(1.6), Inches(0.6),
         font_size=Pt(28), color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, "2026", Inches(11.6), Inches(3.8), Inches(1.6), Inches(0.4),
         font_size=Pt(14), color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# Save
output_path = r"C:\laragon\www\brightdor\BrightDor_Documentation.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
